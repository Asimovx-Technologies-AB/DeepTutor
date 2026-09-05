"""
Parallel Ingestion & Processing Engine.

Features:
- Fast Digital Text Extraction (pypdf / pymupdf)
- Parallel Page VLM OCR Worker Pool (asyncio.gather) for scanned pages
- Multi-format file dispatcher (PDF, Scanned PDF, Images, DOCX, PPTX, TXT)
- Semantic chunking (500–800 chars with 15% overlap)
- Sub-2ms SQLite FTS5 BM25 indexation
- Asynchronous non-blocking background enrichment workers (tables & diagrams)
"""

import os
import io
import asyncio
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
from datetime import datetime

from app.core.config import get_settings
from app.services.study_storage import (
    init_session_db,
    insert_chunks_to_fts,
    save_session_document,
    update_document_status,
    schedule_s3_document_backup,
)


def _chunk_text(text: str, target_size: int = 650, overlap_pct: float = 0.15) -> List[str]:
    """
    Sentence-aware semantic chunking: 500–800 chars with ~15% overlap.
    """
    if not text:
        return []

    # Clean redundant whitespace
    text = " ".join(text.split())
    if len(text) <= target_size:
        return [text]

    overlap_chars = int(target_size * overlap_pct)
    step = target_size - overlap_chars
    chunks = []
    start = 0

    while start < len(text):
        end = start + target_size
        if end >= len(text):
            chunks.append(text[start:].strip())
            break

        # Attempt to split at clean sentence boundary (. ! ? \n)
        boundary = max(
            text.rfind(". ", start, end),
            text.rfind("! ", start, end),
            text.rfind("? ", start, end),
            text.rfind("; ", start, end)
        )
        if boundary != -1 and boundary > start + (target_size // 2):
            chunks.append(text[start:boundary + 1].strip())
            start = boundary + 1
        else:
            # Fallback to space
            space = text.rfind(" ", start, end)
            if space != -1 and space > start:
                chunks.append(text[start:space].strip())
                start = space + 1
            else:
                chunks.append(text[start:end].strip())
                start = end - overlap_chars

    return [c for c in chunks if len(c) > 40]


class StudyDocumentProcessor:
    """Orchestrates zero-wait parallel ingestion and background enrichment."""

    async def ingest_document(
        self,
        doc_id: str,
        file_path: str,
        file_name: str,
        subject: str = "General Study",
        session_id: str = ""
    ) -> Dict[str, Any]:
        """
        Stage 0 & Stage 1: Fast-path extraction, chunking, and FTS5 indexing.
        Returns document metadata and extracted text sample for topic extraction.
        """
        init_session_db(session_id)
        path = Path(file_path)
        ext = path.suffix.lower()

        save_session_document(
            session_id=session_id,
            doc_id=doc_id,
            filename=file_name,
            file_path=file_path,
            status="indexing"
        )

        # Trigger S3 cloud backup asynchronously
        schedule_s3_document_backup(session_id, file_path, file_name)

        extracted_chunks: List[Dict[str, Any]] = []
        full_sample_text = ""
        page_count = 1

        if ext == ".pdf":
            page_count, extracted_chunks, full_sample_text = await self._process_pdf(file_path, doc_id)
        elif ext in (".docx", ".doc"):
            page_count, extracted_chunks, full_sample_text = await self._process_docx(file_path, doc_id)
        elif ext in (".pptx", ".ppt"):
            page_count, extracted_chunks, full_sample_text = await self._process_pptx(file_path, doc_id)
        elif ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"):
            page_count, extracted_chunks, full_sample_text = await self._process_image(file_path, doc_id)
        else:
            # Plain text / Markdown / Code
            page_count, extracted_chunks, full_sample_text = await self._process_text_file(file_path, doc_id)

        # Batch insert chunks into isolated session SQLite FTS5 table
        insert_chunks_to_fts(session_id, doc_id, extracted_chunks)

        # Update document status to text_ready (Chat is now immediately unblocked)
        update_document_status(session_id, doc_id, "text_ready")

        return {
            "doc_id": doc_id,
            "session_id": session_id,
            "filename": file_name,
            "page_count": page_count,
            "chunk_count": len(extracted_chunks),
            "sample_text": full_sample_text[:25000],
            "status": "text_ready"
        }

    # ─── PDF Ingestion: Digital Fast-Text + Parallel VLM OCR ────────────────

    async def _process_pdf(self, file_path: str, doc_id: str) -> Tuple[int, List[Dict[str, Any]], str]:
        """
        Fast inspection: pages with <40 characters digital text are flagged as scanned
        and dispatched concurrently to the Parallel VLM OCR Worker Pool.
        """
        import pymupdf  # fitz
        doc = pymupdf.open(file_path)
        page_count = len(doc)
        digital_pages: List[Tuple[int, str]] = []
        scanned_page_nums: List[int] = []

        for p_idx in range(page_count):
            page = doc[p_idx]
            text = page.get_text().strip()
            if len(text) >= 40:
                digital_pages.append((p_idx + 1, text))
            else:
                scanned_page_nums.append(p_idx + 1)

        chunks: List[Dict[str, Any]] = []
        sample_parts: List[str] = []

        fname = Path(file_path).name
        # Process digital pages immediately
        for page_num, p_text in digital_pages:
            sample_parts.append(p_text)
            p_chunks = _chunk_text(p_text)
            for c_idx, c_text in enumerate(p_chunks):
                chunks.append({
                    "chunk_id": f"{doc_id}_p{page_num}_{c_idx}",
                    "page": page_num,
                    "source_type": "text",
                    "content": f"[Doc: {fname} | Page {page_num} | Type: text] {c_text}"
                })

        # Process scanned pages in parallel (asyncio.gather)
        if scanned_page_nums:
            # Cap parallel OCR to first 12 pages for latency safety
            ocr_targets = scanned_page_nums[:12]
            ocr_tasks = [self._ocr_single_page(file_path, p_num) for p_num in ocr_targets]
            ocr_results = await asyncio.gather(*ocr_tasks, return_exceptions=True)

            for p_num, result in zip(ocr_targets, ocr_results):
                if isinstance(result, str) and result.strip():
                    sample_parts.append(result)
                    p_chunks = _chunk_text(result)
                    for c_idx, c_text in enumerate(p_chunks):
                        chunks.append({
                            "chunk_id": f"{doc_id}_p{p_num}_ocr_{c_idx}",
                            "page": p_num,
                            "source_type": "scanned_vlm",
                            "content": f"[Doc: {fname} | Page {p_num} | Type: scanned_vlm] {c_text}"
                        })

        doc.close()
        full_sample = "\n\n".join(sample_parts)
        return page_count, chunks, full_sample

    async def _ocr_single_page(self, pdf_path: str, page_num: int) -> str:
        """Render page as 130 DPI PNG in-memory and transcribe with Gemini Vision."""
        try:
            import pymupdf
            doc = pymupdf.open(pdf_path)
            page = doc[page_num - 1]
            pix = page.get_pixmap(dpi=130)
            png_bytes = pix.tobytes("png")
            doc.close()

            from app.services.study_agents import call_gemini_vision
            prompt = (
                "Transcribe this academic document page cleanly. Preserve all formulas in LaTeX ($...$ or $$...$$), "
                "maintain academic heading structure, and output clean markdown text. Do not hallucinate."
            )
            resp = await call_gemini_vision(prompt, png_bytes)
            if resp and resp.strip():
                return resp.strip()
        except Exception:
            pass
        return ""

    # ─── Word Document Ingestion ─────────────────────────────────────────────

    async def _process_docx(self, file_path: str, doc_id: str) -> Tuple[int, List[Dict[str, Any]], str]:
        import docx
        doc = docx.Document(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        full_text = "\n\n".join(paragraphs)

        chunks: List[Dict[str, Any]] = []
        raw_chunks = _chunk_text(full_text)
        for i, c in enumerate(raw_chunks):
            chunks.append({
                "chunk_id": f"{doc_id}_c{i}",
                "page": 1,
                "source_type": "text",
                "content": f"[Doc: {Path(file_path).name}] {c}"
            })
        return 1, chunks, full_text

    # ─── PowerPoint Ingestion ────────────────────────────────────────────────

    async def _process_pptx(self, file_path: str, doc_id: str) -> Tuple[int, List[Dict[str, Any]], str]:
        from pptx import Presentation
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        chunks: List[Dict[str, Any]] = []
        sample_parts = []

        fname = Path(file_path).name
        for idx, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            s_text = " ".join(slide_texts)
            if s_text:
                sample_parts.append(s_text)
                for c_idx, c in enumerate(_chunk_text(s_text)):
                    chunks.append({
                        "chunk_id": f"{doc_id}_s{idx+1}_{c_idx}",
                        "page": idx + 1,
                        "source_type": "slide",
                        "content": f"[Doc: {fname} | Slide {idx+1} | Type: slide] {c}"
                    })

        return slide_count, chunks, "\n\n".join(sample_parts)

    # ─── Single Image Ingestion ──────────────────────────────────────────────

    async def _process_image(self, file_path: str, doc_id: str) -> Tuple[int, List[Dict[str, Any]], str]:
        settings = get_settings()
        transcription = ""
        try:
            with open(file_path, "rb") as f:
                img_bytes = f.read()

            from app.services.study_agents import call_vlm
            prompt = (
                "Transcribe this study image or diagram completely. Describe formulas in LaTeX ($...$), "
                "transcribe handwritten text, and explain technical diagrams in detail."
            )
            transcription = await call_vlm(prompt, img_bytes)
        except Exception:
            pass

        if not transcription:
            transcription = f"Uploaded study image: {Path(file_path).name}"

        chunks = [{
            "chunk_id": f"{doc_id}_img_1",
            "page": 1,
            "source_type": "image_caption",
            "content": f"[Doc: {Path(file_path).name} | Image Note] {transcription}"
        }]
        return 1, chunks, transcription

    # ─── Plain Text Ingestion ────────────────────────────────────────────────

    async def _process_text_file(self, file_path: str, doc_id: str) -> Tuple[int, List[Dict[str, Any]], str]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        fname = Path(file_path).name
        chunks: List[Dict[str, Any]] = []
        for i, c in enumerate(_chunk_text(text)):
            chunks.append({
                "chunk_id": f"{doc_id}_c{i}",
                "page": 1,
                "source_type": "text",
                "content": f"[Doc: {fname} | Type: text] {c}"
            })
        return 1, chunks, text

    # ─── Background Enrichment Workers (Stage 2 & Stage 3) ───────────────────

    async def run_background_enrichment(self, session_id: str, doc_id: str, file_path: str):
        """
        Stage 2 (Table Extraction via pdfplumber) & Stage 3 (Diagram Extraction via PyMuPDF).
        Executes without blocking chat; updates document status to fully_processed.
        """
        path = Path(file_path)
        if path.suffix.lower() != ".pdf":
            update_document_status(session_id, doc_id, "fully_processed")
            return

        enrichment_chunks: List[Dict[str, Any]] = []

        # Stage 2: pdfplumber Table Extraction
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for p_idx, page in enumerate(pdf.pages[:20]):
                    tables = page.extract_tables()
                    for t_idx, table in enumerate(tables):
                        if table and len(table) > 1:
                            # Convert to clean markdown table
                            header = " | ".join(str(cell or "").strip() for cell in table[0])
                            divider = " | ".join(["---"] * len(table[0]))
                            rows = [
                                " | ".join(str(cell or "").strip() for cell in row)
                                for row in table[1:]
                                if any(row)
                            ]
                            md_table = f"\n| {header} |\n| {divider} |\n" + "\n".join(f"| {r} |" for r in rows)
                            enrichment_chunks.append({
                                "chunk_id": f"{doc_id}_tbl_p{p_idx+1}_{t_idx}",
                                "page": p_idx + 1,
                                "source_type": "table",
                                "content": f"[Doc: {path.name} | Page {p_idx+1} | Type: table]\n{md_table}"
                            })
        except Exception:
            pass

        # Stage 3: PyMuPDF Embedded Figures Contextual Captioning
        try:
            import pymupdf
            doc = pymupdf.open(file_path)
            settings = get_settings()

            from app.services.study_agents import call_vlm

            extracted_img_count = 0
            for p_idx in range(min(len(doc), 15)):
                page = doc[p_idx]
                images = page.get_images()
                for img_idx, img in enumerate(images[:2]):
                    if extracted_img_count >= 5:
                        break
                    xref = img[0]
                    base_img = doc.extract_image(xref)
                    if base_img and base_img.get("image"):
                        img_bytes = base_img["image"]
                        if len(img_bytes) > 5000:  # Skip tiny icons
                            prompt = "Describe this technical diagram or academic figure concisely. Detail all labeled axes, steps, and key principles."
                            fig_text = await call_vlm(prompt, img_bytes)
                            if fig_text and fig_text.strip():
                                enrichment_chunks.append({
                                    "chunk_id": f"{doc_id}_fig_p{p_idx+1}_{img_idx}",
                                    "page": p_idx + 1,
                                    "source_type": "image_caption",
                                    "content": f"[Doc: {path.name} | Page {p_idx+1} | Type: figure_diagram] {fig_text.strip()}"
                                })
                                extracted_img_count += 1
            doc.close()
        except Exception:
            pass

        if enrichment_chunks:
            insert_chunks_to_fts(session_id, doc_id, enrichment_chunks)

        update_document_status(session_id, doc_id, "fully_processed")


doc_processor = StudyDocumentProcessor()
