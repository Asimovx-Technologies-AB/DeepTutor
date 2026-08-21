"""
Stage 1 — Document Parser
==========================
Multi-engine cascade parser: PyMuPDF → pdfplumber → pypdf → OCR → Docling

For each engine:
  - Extracts raw text per page
  - Detects heading structure (section tree)
  - Preserves formula/LaTeX blocks
  - Hands off to SemanticChunker for 500-1000 word chunks

Architecture (priority order):
  process_document(file_path)
    ├─ PDF:   _try_pymupdf()    ← Primary (fast, layout-aware, formula-preserving)
    │         _try_pdfplumber() ← Fallback 1 (table extraction)
    │         _try_pypdf()      ← Fallback 2 (pure Python)
    │         _ocr_scanned()    ← Fallback 3 (scanned/image PDFs)
    │         _try_docling()    ← Fallback 4 (ML-based, slowest)
    ├─ DOCX:  _try_docx()
    ├─ Image: _try_ocr_image()
    └─ Text/MD/JSON/HTML/CSV/Excel/PPTX: dedicated handlers
"""
from __future__ import annotations

import os
import re
import json
import asyncio
from pathlib import Path
from typing import List, Dict, Optional
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeoutError

from app.core.config import get_settings

settings = get_settings()

# Suppress PyTorch compile errors on Windows
os.environ.setdefault("DOCLING_INFERENCE_COMPILE_TORCH_MODELS", "false")


# ══════════════════════════════════════════════════════════════════════════════
# Helpers
# ══════════════════════════════════════════════════════════════════════════════
_FORMULA_PATTERNS = [
    re.compile(r'\$\$.+?\$\$', re.DOTALL),      # display math $$...$$
    re.compile(r'\$.+?\$'),                       # inline math $...$
    re.compile(r'\\[a-zA-Z]+\{[^}]*\}'),         # LaTeX commands \cmd{...}
    re.compile(r'\\frac\{', re.IGNORECASE),
    re.compile(r'\\[a-zA-Z]{2,}'),               # \alpha, \beta, \theta, etc.
    re.compile(r'(?:[A-Za-z_]+\s*=\s*[-+]?[0-9a-zA-Z_\s\+\-\*/\(\)\^\\\{\}\.]+)', re.MULTILINE), # Algebraic equations
    re.compile(r'(?:[0-9]*[A-Z][a-z]?[0-9]*(?:\([a-z]+\))?\s*\+\s*)+[0-9]*[A-Z][a-z]?[0-9]*(?:\([a-z]+\))?\s*(?:->|→|⇌|=)\s*.+', re.MULTILINE), # Chemical reactions
]


def _detect_formulas(text: str) -> List[str]:
    """Return all formula blocks found in text."""
    formulas = []
    for pattern in _FORMULA_PATTERNS:
        for m in pattern.finditer(text):
            formulas.append(m.group())
    return list(set(formulas))


def _clean_text(text: str) -> str:
    """Normalize whitespace while preserving paragraph structure."""
    text = re.sub(r'\r\n', '\n', text)
    text = re.sub(r'\r', '\n', text)
    text = re.sub(r'[ \t]{2,}', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _is_heading(line: str) -> Optional[tuple[int, str]]:
    """
    Detect heading patterns, returns (level, clean_title) or None.
    Detects: ALL CAPS lines, numbered sections (1.2.3 Title), markdown #
    """
    line = line.strip()
    if not line or len(line) > 120:
        return None
    # Markdown headings
    m = re.match(r'^(#{1,6})\s+(.+)$', line)
    if m:
        return len(m.group(1)), m.group(2).strip()
    # Numbered sections: 1. Introduction / 1.2 Method / 1.2.3 Sub-section
    m = re.match(r'^(\d+(?:\.\d+)*)[.\)]\s+([A-Z].{2,80})$', line)
    if m:
        level = m.group(1).count('.') + 1
        return level, m.group(2).strip()
    # ALL CAPS title (3-8 words)
    if line.isupper() and 2 <= len(line.split()) <= 10 and len(line) <= 80:
        return 1, line.title()
    return None


# ══════════════════════════════════════════════════════════════════════════════
# Hybrid PDF Ingestion Pipeline (PyMuPDF Fast Path + Selective VLM Fallback)
# ══════════════════════════════════════════════════════════════════════════════
def _parse_pdf_hybrid(file_path: str) -> List[Dict]:
    """
    Page-level hybrid document ingestion pipeline:
      1. Evaluates each page: extracts raw text, counts words, and calculates image area coverage.
      2. Traditional Path: Clean digital pages (words >= threshold, low image coverage) are extracted
         instantly via PyMuPDF ($0 cost, sub-millisecond).
      3. VLM Fallback Path: Scanned exam papers, handwritten notes, or diagram-heavy pages
         (words < threshold or image coverage > threshold) are rendered and routed to Gemini Flash VLM.
      4. Per-page disk caching guarantees re-processing the same document incurs 0 duplicate API calls.
      5. Merges pages in original order with extraction_method tag ('traditional' | 'vlm').
    """
    try:
        try:
            import pymupdf as fitz  # Modern PyMuPDF
        except ImportError:
            import fitz              # Legacy fitz fallback

        doc = fitz.open(file_path)
        total_pages = len(doc)
        if total_pages == 0:
            return []

        min_words = getattr(settings, "VLM_MIN_WORDS_THRESHOLD", 50)
        coverage_threshold = getattr(settings, "VLM_IMAGE_COVERAGE_THRESHOLD", 0.70)
        enable_vlm = getattr(settings, "ENABLE_VLM_PARSER", True)
        max_vlm_pages = getattr(settings, "VLM_MAX_PAGES_PER_DOC", 50)

        pages: List[Dict] = []
        vlm_count = 0
        trad_count = 0

        for page_num in range(1, total_pages + 1):
            page = doc[page_num - 1]

            # ── 1. Fast Traditional Text & Word Count ──
            raw_text = page.get_text("text", sort=True) or ""
            cleaned_trad_text = _clean_text(raw_text)
            word_count = len(cleaned_trad_text.split())

            # ── 2. Image Area & Coverage Ratio Calculation ──
            page_area = max(1.0, float(page.rect.width * page.rect.height))
            image_area = 0.0
            image_list = page.get_images() or []

            for img in image_list:
                xref = img[0]
                for rect in page.get_image_rects(xref):
                    image_area += float(rect.width * rect.height)

            image_coverage = min(1.0, image_area / page_area)

            # ── 3. Page-Level Scan Detection ──
            is_blank = (word_count < 10 and image_coverage < 0.05 and len(image_list) == 0)
            needs_vlm = False

            if enable_vlm and not is_blank and vlm_count < max_vlm_pages:
                # Flag A: Scanned page or low text density with images
                if word_count < min_words and (image_coverage > 0.05 or len(image_list) > 0):
                    needs_vlm = True
                # Flag B: Dominated by single/multiple large images (> 70% coverage)
                elif image_coverage >= coverage_threshold:
                    needs_vlm = True
                # Flag C: Handwritten / question diagram page with very few digital words
                elif word_count < 15 and len(image_list) > 0:
                    needs_vlm = True

            # ── 4. Execution (Traditional vs VLM) ──
            final_text = ""
            extraction_method = "traditional"

            if needs_vlm:
                try:
                    pix = page.get_pixmap(dpi=180)
                    img_bytes = pix.tobytes("png")
                    vlm_prompt = (
                        "Extract all readable text, accurately format all mathematical and scientific formulas in LaTeX notation ($...$ for inline, $$...$$ for display equations), "
                        "provide detailed pedagogical descriptions of any diagrams/figures/charts, and reconstruct any tables as clean Markdown tables."
                    )
                    vlm_result = _try_gemini_vlm(img_bytes, prompt=vlm_prompt)
                    if vlm_result and len(vlm_result.strip()) >= 15:
                        final_text = vlm_result
                        extraction_method = "vlm"
                        vlm_count += 1
                    else:
                        # Fallback to traditional text if VLM failed or returned empty
                        final_text = cleaned_trad_text if cleaned_trad_text else "[Visual Content: Image/Diagram with low text clarity]"
                        trad_count += 1
                except Exception as e:
                    print(f"[HYBRID PARSER WARN] Page {page_num} VLM fallback to traditional: {e}")
                    final_text = cleaned_trad_text if cleaned_trad_text else "[Visual Content: Image/Diagram]"
                    trad_count += 1
            else:
                final_text = cleaned_trad_text
                trad_count += 1

            # ── 5. Clean Markdown, Formulas & Section Hierarchy ──
            final_clean = _clean_text(final_text)
            headings = []
            for line in final_clean.split('\n'):
                h = _is_heading(line)
                if h:
                    headings.append({"level": h[0], "title": h[1]})

            formulas = _detect_formulas(final_clean)

            pages.append({
                "page": page_num,
                "text": final_clean,
                "headings": headings,
                "formulas": formulas,
                "char_count": len(final_clean),
                "word_count": len(final_clean.split()),
                "extraction_method": extraction_method,
                "image_coverage": round(image_coverage, 3),
            })

        doc.close()
        print(f"[HYBRID PARSER] '{Path(file_path).name}' ({len(pages)} pages) -> {trad_count} traditional, {vlm_count} VLM")
        return pages

    except Exception as e:
        print(f"[HYBRID PARSER ERROR] PyMuPDF error on {file_path}: {e}")
        return []



# ══════════════════════════════════════════════════════════════════════════════
# pdfplumber Parser (Fallback 1)
# ══════════════════════════════════════════════════════════════════════════════
def _try_pdfplumber(file_path: str) -> List[Dict]:
    """Fallback 1: pdfplumber — good for table extraction."""
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                text = _clean_text(text)
                if len(text) < 20:
                    continue
                # Extract tables as markdown
                table_blocks = []
                for table in page.extract_tables() or []:
                    if not table:
                        continue
                    lines = []
                    for i, row in enumerate(table):
                        cells = [str(c or "").replace("\n", " ").strip() for c in row]
                        lines.append("| " + " | ".join(cells) + " |")
                        if i == 0:
                            lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                    table_blocks.append("\n".join(lines))

                combined = text
                if table_blocks:
                    combined = text + "\n\n" + "\n\n".join(table_blocks)

                pages.append({
                    "page": page_num,
                    "text": combined,
                    "headings": [],
                    "formulas": _detect_formulas(text),
                    "char_count": len(combined),
                })
        return pages
    except ImportError:
        return []
    except Exception as e:
        print(f"[PARSER pdfplumber] Error: {e}")
        return []


# ══════════════════════════════════════════════════════════════════════════════
# pypdf Parser (Fallback 2)
# ══════════════════════════════════════════════════════════════════════════════
def _try_pypdf(file_path: str) -> List[Dict]:
    """Fallback 2: pypdf pure-Python parser."""
    try:
        import pypdf
        pages = []
        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                text = _clean_text(text)
                if len(text) < 20:
                    continue
                pages.append({
                    "page": page_num,
                    "text": text,
                    "headings": [],
                    "formulas": _detect_formulas(text),
                    "char_count": len(text),
                })
        return pages
    except ImportError:
        return []
    except Exception as e:
        print(f"[PARSER pypdf] Error: {e}")
        return []


# ══════════════════════════════════════════════════════════════════════════════
# OCR for Scanned PDFs (Fallback 3)
# ══════════════════════════════════════════════════════════════════════════════
_easyocr_reader = None

def _get_easyocr():
    global _easyocr_reader
    if _easyocr_reader is not None:
        return _easyocr_reader
    try:
        import easyocr, torch
        use_gpu = torch.cuda.is_available()
        print(f"[OCR] Initializing EasyOCR (GPU={use_gpu})...")
        _easyocr_reader = easyocr.Reader(['en'], gpu=use_gpu, verbose=False)
    except Exception as e:
        print(f"[OCR] EasyOCR init failed: {e}")
    return _easyocr_reader


def _try_gemini_vlm(image_input, prompt: Optional[str] = None) -> str:
    """Run Google Gemini 2.0 / 1.5 Flash VLM to transcribe an image or document page."""
    if not getattr(settings, "ENABLE_VLM_PARSER", True):
        return ""
    try:
        from app.rag.gemini_client import gemini_client
        api_key = os.environ.get("GEMINI_API_KEY") or getattr(settings, "GEMINI_API_KEY", "")
        if not api_key or len(api_key.strip()) < 5:
            return ""
        
        target_vlm_model = getattr(settings, "GEMINI_VLM_MODEL", "gemini-2.0-flash") or "gemini-2.0-flash"
        text = gemini_client.sync_transcribe_image_vlm(
            image_input,
            prompt=prompt,
            model=target_vlm_model
        )
        if text and len(text.strip()) >= 10:
            return text.strip()
    except Exception as e:
        print(f"[VLM WARN] Gemini VLM transcription error: {e}")
    return ""


def _ocr_page_image(pil_image) -> str:
    """Run Gemini Flash VLM → EasyOCR → PyTesseract fallback on a PIL image."""
    # 1. Primary: Gemini Flash VLM
    vlm_text = _try_gemini_vlm(pil_image)
    if vlm_text:
        return vlm_text

    # 2. Secondary: EasyOCR
    reader = _get_easyocr()
    if reader:
        try:
            results = reader.readtext(pil_image, detail=0)
            text = " ".join(results).strip()
            if len(text) >= 10:
                return text
        except Exception:
            pass

    # 3. Tertiary: PyTesseract
    try:
        import pytesseract
        text = pytesseract.image_to_string(pil_image).strip()
        if len(text) >= 10:
            return text
    except Exception:
        pass
    return ""


def _ocr_scanned_pdf(file_path: str) -> List[Dict]:
    """Render scanned PDF pages as images and transcribe them using Gemini Flash VLM / OCR."""
    pages = []
    max_pages = getattr(settings, "VLM_MAX_PAGES_PER_DOC", 50) or 50
    try:
        try:
            import pymupdf as fitz
        except ImportError:
            import fitz
        doc = fitz.open(file_path)
        total_pages = min(len(doc), max_pages)
        print(f"[VLM/OCR] Processing {total_pages} scanned PDF pages from {Path(file_path).name} (VLM={getattr(settings, 'ENABLE_VLM_PARSER', True)})...")

        for page_num in range(1, total_pages + 1):
            page = doc[page_num - 1]
            pix = page.get_pixmap(dpi=180)
            try:
                from PIL import Image
                import io
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                text = _ocr_page_image(img)
            except Exception:
                text = ""

            if text and len(text) >= 15:
                cleaned = _clean_text(text)
                headings = []
                for line in cleaned.split('\n'):
                    res = _is_heading(line)
                    if res:
                        headings.append({"level": res[0], "title": res[1]})

                pages.append({
                    "page": page_num,
                    "text": cleaned,
                    "headings": headings,
                    "formulas": _detect_formulas(cleaned),
                    "char_count": len(cleaned),
                })
        doc.close()
    except Exception as e:
        print(f"[VLM/OCR] Scanned PDF processing error: {e}")
    return pages



# ══════════════════════════════════════════════════════════════════════════════
# Docling Parser (Fallback 4 — ML-based, slowest)
# ══════════════════════════════════════════════════════════════════════════════
_docling_converter = None
_docling_available = None


def _get_docling():
    global _docling_converter, _docling_available
    if _docling_available is False:
        return None
    if _docling_converter is not None:
        return _docling_converter
    try:
        import torch
        from docling.datamodel.base_models import InputFormat
        from docling.datamodel.pipeline_options import PdfPipelineOptions
        from docling.document_converter import DocumentConverter, PdfFormatOption

        opts = PdfPipelineOptions()
        opts.do_ocr = settings.DOCLING_ENABLE_OCR
        if torch.cuda.is_available():
            from docling.datamodel.pipeline_options import AcceleratorOptions
            opts.accelerator_options = AcceleratorOptions(num_threads=8, device="cuda:0")

        _docling_converter = DocumentConverter(
            format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=opts)}
        )
        _docling_available = True
        return _docling_converter
    except Exception as e:
        print(f"[DOCLING] Init failed: {e}")
        _docling_available = False
        return None


def _try_docling(file_path: str) -> List[Dict]:
    if not settings.ENABLE_DOCLING:
        return []
    def _run():
        converter = _get_docling()
        if not converter:
            return []
        result = converter.convert(file_path)
        md = result.document.export_to_markdown()
        if not md:
            return []
        return [{"page": 1, "text": _clean_text(md), "headings": [], "formulas": _detect_formulas(md), "char_count": len(md)}]
    try:
        with ThreadPoolExecutor(max_workers=1) as ex:
            return ex.submit(_run).result(timeout=settings.DOCLING_TIMEOUT_SECONDS)
    except Exception:
        return []


# ══════════════════════════════════════════════════════════════════════════════
# DocumentParser — Public API
# ══════════════════════════════════════════════════════════════════════════════
class DocumentParser:
    """
    Multi-engine cascade document parser.
    Returns a list of page dicts: {page, text, headings, formulas, char_count}
    """

    def parse_pdf(self, file_path: str) -> List[Dict]:
        """PDF Hybrid Pipeline: PyMuPDF Fast Path with selective Gemini Flash VLM per-page fallback."""
        # 1. Primary: Hybrid Page-by-Page Ingestion
        pages = _parse_pdf_hybrid(file_path)
        if pages and any(p.get("char_count", 0) > 20 for p in pages):
            return pages

        # 2. pdfplumber (tables fallback)
        pages = _try_pdfplumber(file_path)
        if pages and sum(p.get("char_count", 0) for p in pages) > 200:
            print(f"[PARSER] pdfplumber: {len(pages)} pages")
            for p in pages:
                p.setdefault("extraction_method", "traditional")
            return pages

        # 3. pypdf (pure python fallback)
        pages = _try_pypdf(file_path)
        if pages and sum(p.get("char_count", 0) for p in pages) > 200:
            print(f"[PARSER] pypdf: {len(pages)} pages")
            for p in pages:
                p.setdefault("extraction_method", "traditional")
            return pages

        # 4. OCR / Scanned PDF fallback
        print(f"[PARSER] Running fallback OCR/VLM on scanned PDF: {Path(file_path).name}")
        pages = _ocr_scanned_pdf(file_path)
        if pages:
            return pages

        # 5. Docling (ML-based)
        return _try_docling(file_path)


    def parse_docx(self, file_path: str) -> List[Dict]:
        try:
            import docx
            doc = docx.Document(file_path)
            source = Path(file_path).name
            blocks, section = [], ""
            for p in doc.paragraphs:
                t = p.text.strip()
                if not t:
                    continue
                if p.style and p.style.name.startswith("Heading"):
                    section = t
                    blocks.append(f"\n# {t}\n")
                else:
                    blocks.append(t)
            for table in doc.tables:
                rows = []
                for i, row in enumerate(table.rows):
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                blocks.append("\n" + "\n".join(rows))
            text = _clean_text("\n".join(blocks))
            return [{"page": 1, "text": text, "headings": [], "formulas": _detect_formulas(text), "char_count": len(text)}] if text else []
        except Exception:
            return self.parse_txt(file_path)

    def parse_image(self, file_path: str) -> List[Dict]:
        """Parse standalone image/diagram via Gemini Flash VLM (with OCR/Docling fallback)."""
        source = Path(file_path).name

        # 1. Primary: Gemini Flash VLM
        text = _try_gemini_vlm(file_path)

        # 2. Fallback: OCR
        if not text:
            text = _ocr_page_image(file_path)

        # 3. Fallback: Docling
        if not text and getattr(settings, "ENABLE_DOCLING", False):
            docling_pages = _try_docling(file_path)
            if docling_pages:
                return docling_pages

        if text:
            full = text if text.startswith("#") else f"# Image / Diagram Document: {source}\n\n{text}"
            headings = []
            for line in full.split('\n'):
                res = _is_heading(line)
                if res:
                    headings.append({"level": res[0], "title": res[1]})

            return [{
                "page": 1,
                "text": full,
                "headings": headings,
                "formulas": _detect_formulas(full),
                "char_count": len(full),
            }]
        return []


    def parse_txt(self, file_path: str) -> List[Dict]:
        try:
            text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
            text = _clean_text(text)
            if not text:
                return []
            return [{"page": 1, "text": text, "headings": [], "formulas": _detect_formulas(text), "char_count": len(text)}]
        except Exception:
            return []

    def parse_csv(self, file_path: str) -> List[Dict]:
        try:
            import csv
            source = Path(file_path).name
            rows = []
            with open(file_path, encoding="utf-8", errors="ignore") as f:
                for row in csv.reader(f):
                    if row:
                        rows.append(row)
            if not rows:
                return []
            header = rows[0]
            batch_size = 40
            pages = []
            for i in range(0, max(1, len(rows[1:])), batch_size):
                batch = rows[1:][i:i + batch_size]
                lines = [
                    f"# CSV: {source} (rows {i+1}–{i+len(batch)})",
                    "| " + " | ".join(str(h) for h in header) + " |",
                    "| " + " | ".join(["---"] * len(header)) + " |",
                ]
                for row in batch:
                    lines.append("| " + " | ".join(str(v).replace("\n", " ") for v in row) + " |")
                text = "\n".join(lines)
                pages.append({"page": i // batch_size + 1, "text": text, "headings": [], "formulas": [], "char_count": len(text)})
            return pages
        except Exception:
            return self.parse_txt(file_path)

    def parse_excel(self, file_path: str) -> List[Dict]:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            pages = []
            for page_num, sheet_name in enumerate(wb.sheetnames, 1):
                sheet = wb[sheet_name]
                rows = [[str(c or "").strip() for c in r] for r in sheet.iter_rows(values_only=True) if any(c is not None for c in r)]
                if not rows:
                    continue
                lines = [f"# Excel Sheet: {sheet_name}", "| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * len(rows[0])) + " |"]
                for r in rows[1:]:
                    lines.append("| " + " | ".join(r) + " |")
                text = "\n".join(lines)
                pages.append({"page": page_num, "text": text, "headings": [{"level": 1, "title": f"Sheet: {sheet_name}"}], "formulas": [], "char_count": len(text)})
            return pages
        except Exception:
            return self.parse_txt(file_path)

    def parse_pptx(self, file_path: str) -> List[Dict]:
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            pages = []
            for slide_num, slide in enumerate(prs.slides, 1):
                texts, title = [], ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if t:
                            if shape == slide.shapes.title:
                                title = t
                            else:
                                texts.append(t)
                text = f"# Slide {slide_num}: {title}\n\n" + "\n\n".join(texts)
                if len(text.strip()) >= 10:
                    pages.append({"page": slide_num, "text": text.strip(), "headings": [{"level": 1, "title": f"Slide {slide_num}: {title}"}], "formulas": [], "char_count": len(text)})
            return pages
        except Exception:
            return self.parse_txt(file_path)

    def parse_html(self, file_path: str) -> List[Dict]:
        try:
            from bs4 import BeautifulSoup
            raw = Path(file_path).read_text(encoding="utf-8", errors="ignore")
            soup = BeautifulSoup(raw, "html.parser")
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()
            text = _clean_text(soup.get_text(separator="\n"))
            return [{"page": 1, "text": text, "headings": [], "formulas": _detect_formulas(text), "char_count": len(text)}] if text else []
        except Exception:
            return self.parse_txt(file_path)

    def parse_json(self, file_path: str) -> List[Dict]:
        try:
            data = json.loads(Path(file_path).read_text(encoding="utf-8", errors="ignore"))
            text = json.dumps(data, indent=2)
            return [{"page": 1, "text": text, "headings": [], "formulas": [], "char_count": len(text)}]
        except Exception:
            return self.parse_txt(file_path)

    def parse(self, file_path: str) -> List[Dict]:
        """Dispatch to the correct parser based on file extension."""
        ext = Path(file_path).suffix.lower()
        dispatch = {
            ".pdf":  self.parse_pdf,
            ".docx": self.parse_docx,
            ".doc":  self.parse_docx,
            ".xlsx": self.parse_excel,
            ".xls":  self.parse_excel,
            ".pptx": self.parse_pptx,
            ".ppt":  self.parse_pptx,
            ".csv":  self.parse_csv,
            ".html": self.parse_html,
            ".htm":  self.parse_html,
            ".json": self.parse_json,
            ".txt":  self.parse_txt,
            ".md":   self.parse_txt,
            ".rst":  self.parse_txt,
            ".log":  self.parse_txt,
            ".py":   self.parse_txt,
            ".js":   self.parse_txt,
            ".ts":   self.parse_txt,
            ".png":  self.parse_image,
            ".jpg":  self.parse_image,
            ".jpeg": self.parse_image,
            ".webp": self.parse_image,
            ".bmp":  self.parse_image,
            ".tiff": self.parse_image,
            ".tif":  self.parse_image,
        }
        handler = dispatch.get(ext, self.parse_txt)
        return handler(file_path)


# Singleton
document_parser = DocumentParser()
